import io
import pypdf
from docx import Document
from pptx import Presentation

def extract_text_from_file(uploaded_file):
    """
    Extracts text from an uploaded file (PDF, DOCX, PPTX).
    Returns the extracted text as a string.
    """
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    if file_type == 'pdf':
        return extract_text_from_pdf(uploaded_file)
    elif file_type == 'docx':
        return extract_text_from_docx(uploaded_file)
    elif file_type == 'pptx':
        return extract_text_from_pptx(uploaded_file)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

def extract_text_from_pdf(file):
    text = ""
    try:
        pdf_reader = pypdf.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        raise Exception(f"Error reading PDF: {e}")
    return text

def extract_text_from_docx(file):
    text = ""
    try:
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        raise Exception(f"Error reading DOCX: {e}")
    return text

def extract_text_from_pptx(file):
    text = ""
    try:
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise Exception(f"Error reading PPTX: {e}")
    return text
